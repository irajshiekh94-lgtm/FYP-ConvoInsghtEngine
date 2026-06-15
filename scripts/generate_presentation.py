#!/usr/bin/env python3
"""Generate ConvoInsight FYP pitch deck (PowerPoint)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "ConvoInsight_Presentation.pptx"
LOGO = ROOT / "frontend" / "assets" / "images" / "ConvoInsightLogo.png"

# Brand palette (matches app theme)
BRAND_DARK = RGBColor(11, 20, 26)
BRAND_GREEN = RGBColor(0, 128, 105)
BRAND_GREEN_BRIGHT = RGBColor(0, 168, 132)
BG_LIGHT = RGBColor(248, 250, 252)
BG_WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(17, 27, 33)
TEXT_MUTED = RGBColor(102, 119, 129)
TEXT_ON_GREEN = RGBColor(255, 255, 255)
ACCENT_LINE = RGBColor(0, 168, 132)
CARD_FILL = RGBColor(255, 255, 255)
CARD_BORDER = RGBColor(225, 232, 237)
SIDEBAR_WIDTH = Inches(0.62)


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self._slide_num = 0

    def _blank(self):
        self._slide_num += 1
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    @staticmethod
    def _solid(shape, color: RGBColor) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    @staticmethod
    def _set_bg(slide, color: RGBColor) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _footer(self, slide, label: str = "ConvoInsight Engine") -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.12), Inches(10), Inches(0.38)
        )
        self._solid(bar, BRAND_DARK)
        box = slide.shapes.add_textbox(Inches(0.55), Inches(7.16), Inches(5), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(160, 180, 190)
        p.font.bold = True

        num_box = slide.shapes.add_textbox(Inches(9.1), Inches(7.16), Inches(0.6), Inches(0.3))
        np = num_box.text_frame.paragraphs[0]
        np.text = f"{self._slide_num:02d}"
        np.font.size = Pt(9)
        np.font.color.rgb = BRAND_GREEN_BRIGHT
        np.alignment = PP_ALIGN.RIGHT

    def _sidebar(self, slide) -> None:
        panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SIDEBAR_WIDTH, Inches(7.12)
        )
        self._solid(panel, BRAND_GREEN)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, SIDEBAR_WIDTH, Inches(0), Inches(0.06), Inches(7.12)
        )
        self._solid(accent, BRAND_GREEN_BRIGHT)

    def _title_block(self, slide, title: str, *, x: float = 0.95, y: float = 0.42) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(8.5), Inches(0.75))
        p = box.text_frame.paragraphs[0]
        p.text = title.upper()
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.72), Inches(1.35), Inches(0.07)
        )
        self._solid(line, BRAND_GREEN_BRIGHT)

    def _deco_circles(self, slide) -> None:
        specs = [
            (8.2, 0.4, 2.4, RGBColor(0, 168, 132)),
            (7.4, 5.0, 3.0, RGBColor(0, 100, 82)),
            (-0.8, 4.8, 2.2, RGBColor(0, 140, 115)),
        ]
        for x, y, size, color in specs:
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.fill.transparency = 0.88
            circle.line.fill.background()

    def add_title_slide(self, title: str, subtitle: str) -> None:
        slide = self._blank()
        self._set_bg(slide, BRAND_DARK)
        self._deco_circles(slide)

        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.1)
        )
        self._solid(top_bar, BRAND_GREEN_BRIGHT)

        if LOGO.exists():
            slide.shapes.add_picture(str(LOGO), Inches(3.85), Inches(0.85), height=Inches(1.35))

        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.55), Inches(2.35), Inches(2.9), Inches(0.42)
        )
        self._solid(badge, BRAND_GREEN)
        badge.adjustments[0] = 0.5
        badge_tf = badge.text_frame
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = badge_tf.paragraphs[0]
        bp.text = "FINAL YEAR PROJECT"
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = TEXT_ON_GREEN
        bp.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.95), Inches(8.5), Inches(1.1))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(46)
        tp.font.bold = True
        tp.font.color.rgb = TEXT_ON_GREEN
        tp.alignment = PP_ALIGN.CENTER

        sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.05), Inches(7.6), Inches(1.2))
        stf = sub_box.text_frame
        stf.word_wrap = True
        for i, line in enumerate(subtitle.split("\n")):
            p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
            p.text = line
            p.font.size = Pt(19)
            p.font.color.rgb = RGBColor(180, 210, 200)
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(4)

        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(5.55), Inches(4.4), Inches(0.55)
        )
        pill.fill.background()
        pill.line.color.rgb = BRAND_GREEN_BRIGHT
        pill.line.width = Pt(1.5)
        pill.adjustments[0] = 0.5
        ptf = pill.text_frame
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]
        pp.text = "WhatsApp  ·  AI  ·  Mobile  ·  Business Intelligence"
        pp.font.size = Pt(12)
        pp.font.color.rgb = BRAND_GREEN_BRIGHT
        pp.alignment = PP_ALIGN.CENTER

        self._footer(slide, "ConvoInsight Engine  |  Pitch Deck")

    def add_content_slide(self, title: str, bullets: list[str]) -> None:
        slide = self._blank()
        self._set_bg(slide, BG_LIGHT)
        self._sidebar(slide)
        self._title_block(slide, title)
        self._footer(slide)

        y = 1.55
        for bullet in bullets:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.95),
                Inches(y),
                Inches(8.75),
                Inches(0.62),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_FILL
            card.line.color.rgb = CARD_BORDER
            card.line.width = Pt(0.75)
            card.adjustments[0] = 0.08

            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y), Inches(0.08), Inches(0.62)
            )
            self._solid(stripe, BRAND_GREEN_BRIGHT)

            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1.12), Inches(y + 0.21), Inches(0.12), Inches(0.12)
            )
            self._solid(dot, BRAND_GREEN)

            box = slide.shapes.add_textbox(Inches(1.32), Inches(y + 0.1), Inches(8.2), Inches(0.45))
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(14.5)
            p.font.color.rgb = TEXT_DARK
            y += 0.72

    def add_flow_slide(self, title: str, bullets: list[str]) -> None:
        slide = self._blank()
        self._set_bg(slide, BG_LIGHT)
        self._sidebar(slide)
        self._title_block(slide, title)
        self._footer(slide)

        for i, bullet in enumerate(bullets):
            col = i % 2
            row = i // 2
            x = 0.95 + col * 4.45
            y = 1.55 + row * 1.75

            step = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.48), Inches(0.48)
            )
            self._solid(step, BRAND_GREEN if i % 2 == 0 else BRAND_GREEN_BRIGHT)
            stf = step.text_frame
            stf.vertical_anchor = MSO_ANCHOR.MIDDLE
            sp = stf.paragraphs[0]
            sp.text = str(i + 1)
            sp.font.size = Pt(14)
            sp.font.bold = True
            sp.font.color.rgb = TEXT_ON_GREEN
            sp.alignment = PP_ALIGN.CENTER

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.58),
                Inches(y - 0.04),
                Inches(3.75),
                Inches(0.95),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_FILL
            card.line.color.rgb = CARD_BORDER
            card.line.width = Pt(0.75)
            card.adjustments[0] = 0.1

            box = slide.shapes.add_textbox(Inches(x + 0.72), Inches(y + 0.08), Inches(3.45), Inches(0.75))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            text = bullet
            if text[0].isdigit() and ". " in text[:4]:
                text = text.split(". ", 1)[1]
            p.text = text
            p.font.size = Pt(12.5)
            p.font.color.rgb = TEXT_DARK

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left: list[str],
        right_title: str,
        right: list[str],
    ) -> None:
        slide = self._blank()
        self._set_bg(slide, BG_LIGHT)
        self._sidebar(slide)
        self._title_block(slide, title)
        self._footer(slide)

        for col_x, hdr_title, items, accent in [
            (0.95, left_title, left, BRAND_GREEN),
            (5.15, right_title, right, BRAND_GREEN_BRIGHT),
        ]:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(col_x),
                Inches(1.45),
                Inches(4.05),
                Inches(5.35),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_FILL
            card.line.color.rgb = CARD_BORDER
            card.line.width = Pt(0.75)
            card.adjustments[0] = 0.04

            hdr = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(col_x),
                Inches(1.45),
                Inches(4.05),
                Inches(0.62),
            )
            self._solid(hdr, accent)

            hbox = slide.shapes.add_textbox(
                Inches(col_x + 0.2), Inches(1.58), Inches(3.65), Inches(0.4)
            )
            hp = hbox.text_frame.paragraphs[0]
            hp.text = hdr_title
            hp.font.size = Pt(16)
            hp.font.bold = True
            hp.font.color.rgb = TEXT_ON_GREEN

            y = 2.25
            for item in items:
                row = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(col_x + 0.18),
                    Inches(y),
                    Inches(3.68),
                    Inches(0.48),
                )
                row.fill.solid()
                row.fill.fore_color.rgb = BG_LIGHT
                row.line.fill.background()
                row.adjustments[0] = 0.15

                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(col_x + 0.28),
                    Inches(y + 0.16),
                    Inches(0.1),
                    Inches(0.1),
                )
                self._solid(dot, accent)

                box = slide.shapes.add_textbox(
                    Inches(col_x + 0.48), Inches(y + 0.08), Inches(3.2), Inches(0.35)
                )
                p = box.text_frame.paragraphs[0]
                p.text = item
                p.font.size = Pt(12.5)
                p.font.color.rgb = TEXT_DARK
                y += 0.58

    def add_pipeline_slide(self, title: str, bullets: list[str]) -> None:
        slide = self._blank()
        self._set_bg(slide, BG_LIGHT)
        self._sidebar(slide)
        self._title_block(slide, title)
        self._footer(slide)

        for i, bullet in enumerate(bullets):
            col = i % 3
            row = i // 3
            x = 0.95 + col * 3.05
            y = 1.55 + row * 1.55

            num = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.42), Inches(0.42)
            )
            self._solid(num, BRAND_GREEN if row == 0 else BRAND_GREEN_BRIGHT)
            ntf = num.text_frame
            ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
            np = ntf.paragraphs[0]
            np.text = str(i + 1)
            np.font.size = Pt(12)
            np.font.bold = True
            np.font.color.rgb = TEXT_ON_GREEN
            np.alignment = PP_ALIGN.CENTER

            box = slide.shapes.add_textbox(Inches(x + 0.02), Inches(y + 0.48), Inches(2.85), Inches(0.95))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            text = bullet.split(". ", 1)[1] if ". " in bullet else bullet
            p.text = text
            p.font.size = Pt(11.5)
            p.font.color.rgb = TEXT_DARK
            p.alignment = PP_ALIGN.CENTER

    def add_closing_slide(self, title: str, bullets: list[str]) -> None:
        slide = self._blank()
        self._set_bg(slide, BRAND_DARK)
        self._deco_circles(slide)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(8.4), Inches(0.8))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = title.upper()
        tp.font.size = Pt(30)
        tp.font.bold = True
        tp.font.color.rgb = TEXT_ON_GREEN
        tp.alignment = PP_ALIGN.CENTER

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(1.42), Inches(1.6), Inches(0.07)
        )
        self._solid(line, BRAND_GREEN_BRIGHT)

        y = 1.85
        for bullet in bullets[:-1]:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.1),
                Inches(y),
                Inches(7.8),
                Inches(0.58),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 35, 42)
            card.line.color.rgb = RGBColor(0, 100, 82)
            card.line.width = Pt(0.75)
            card.adjustments[0] = 0.1

            box = slide.shapes.add_textbox(Inches(1.35), Inches(y + 0.12), Inches(7.3), Inches(0.4))
            p = box.text_frame.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(220, 235, 230)
            y += 0.72

        thanks = slide.shapes.add_textbox(Inches(1.5), Inches(6.35), Inches(7), Inches(0.6))
        thp = thanks.text_frame.paragraphs[0]
        thp.text = bullets[-1]
        thp.font.size = Pt(24)
        thp.font.bold = True
        thp.font.color.rgb = BRAND_GREEN_BRIGHT
        thp.alignment = PP_ALIGN.CENTER

        self._footer(slide, "ConvoInsight Engine  |  Thank You")


def build() -> Path:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    deck = Deck()

    deck.add_title_slide(
        "ConvoInsight Engine",
        "AI-Powered WhatsApp Chat Analysis\nFinal Year Project",
    )

    deck.add_content_slide(
        "Problem Statement",
        [
            "WhatsApp group and business chats generate hundreds of messages daily.",
            "Important tasks, deadlines, and decisions get buried in casual conversation.",
            "Manually reading long chat exports is slow, error-prone, and not scalable.",
            "Users need a way to turn raw chat history into actionable business insights.",
        ],
    )

    deck.add_content_slide(
        "Our Solution — ConvoInsight",
        [
            "A mobile app that analyzes exported WhatsApp chats using AI and NLP.",
            "Upload a chat export → automated pipeline extracts insights in minutes.",
            "Surfaces priorities, action items, summaries, sentiment, and topics.",
            "Designed for students, teams, and professionals managing chat-heavy workflows.",
        ],
    )

    deck.add_content_slide(
        "Key Features",
        [
            "📱 Email OTP login & signup — secure, passwordless authentication",
            "📂 Import WhatsApp chat exports (.txt) directly from the phone",
            "⚡ AI executive summary — 24-hour and full-conversation summaries (Llama 3.2)",
            "🎯 Priority dashboard — urgent, moderate, and low-priority messages",
            "✅ Action item extraction — tasks, meetings, and follow-ups",
            "🔔 Urgent notifications — local alerts for high-priority messages",
            "🌓 Light / dark theme with modern WhatsApp-inspired UI",
        ],
    )

    deck.add_flow_slide(
        "How It Works — User Flow",
        [
            "1. Sign up / Log in with email → receive OTP verification code",
            "2. Export chat from WhatsApp (without media) → upload in app",
            "3. Backend parses, cleans, and analyzes messages in background",
            "4. View Dashboard tab: summaries, priorities, actions, analytics",
            "5. View Conversation tab: original messages with read/unread tracking",
            "6. Re-run ⚡ summary anytime for latest 24-hour insights",
        ],
    )

    deck.add_two_column_slide(
        "System Architecture",
        "Frontend (Mobile)",
        [
            "React Native + Expo",
            "TypeScript",
            "React Navigation",
            "TanStack React Query",
            "AsyncStorage",
            "Expo Notifications",
        ],
        "Backend (API Server)",
        [
            "Python FastAPI",
            "Uvicorn ASGI server",
            "REST API (/api/*)",
            "Background job pipeline",
            "MongoDB Atlas storage",
            "Gmail SMTP (OTP emails)",
        ],
    )

    deck.add_pipeline_slide(
        "Analysis Pipeline (9 Steps)",
        [
            "1. Parse WhatsApp export text",
            "2. Clean & normalize messages",
            "3. Structure into standardized message objects",
            "4. Cluster messages by topic (ML clustering)",
            "5. Summarize clusters (Llama 3.2 via Ollama)",
            "6. Classify intent (requests, questions, decisions…)",
            "7. Assign priority levels (urgent / moderate / low)",
            "8. Extract action items & meetings",
            "9. Enrich with entities, topics, sentiment & analytics",
        ],
    )

    deck.add_two_column_slide(
        "AI & Machine Learning Stack",
        "Core ML / NLP",
        [
            "Meta Llama 3.2 (Ollama — local)",
            "OpenAI Whisper (speech-to-text)",
            "scikit-learn (topic clustering)",
            "Transformers / PyTorch",
            "Intent classification pipeline",
            "Priority & action extractors",
        ],
        "Supporting Services",
        [
            "Google Generative AI (legacy intent)",
            "Rule-based summarization fallback",
            "Text normalization & cleaning",
            "Entity & topic extraction",
            "Sentiment analysis",
            "Insight enrichment layer",
        ],
    )

    deck.add_content_slide(
        "Backend Technologies",
        [
            "FastAPI — high-performance Python REST framework",
            "MongoDB Atlas — users, messages, analysis jobs (persistent storage)",
            "Pydantic — request/response validation & schemas",
            "Gmail SMTP — sends OTP codes to any user's email",
            "Ollama — local LLM inference (llama3.2 model)",
            "pytest + httpx — automated API testing",
        ],
    )

    deck.add_content_slide(
        "Frontend & Mobile Technologies",
        [
            "React Native 0.76 — cross-platform iOS & Android",
            "Expo SDK 52 — build, notifications, document picker",
            "TypeScript — type-safe client code",
            "React Navigation 7 — stack & tab navigation",
            "Custom design system — themed components, light/dark mode",
            "Lazy-loaded screens — faster app startup",
        ],
    )

    deck.add_closing_slide(
        "Conclusion & Future Work",
        [
            "ConvoInsight transforms unstructured WhatsApp chats into structured business intelligence.",
            "Combines modern mobile UX with a robust AI/NLP backend pipeline.",
            "Future: real-time WhatsApp integration, multi-language support, team dashboards.",
            "Future: cloud LLM scaling (Groq/Together), advanced analytics & export reports.",
            "Thank you — Questions?",
        ],
    )

    deck.prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Created: {path}")
